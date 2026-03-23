from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
import math
import re
from pathlib import Path
from typing import Any

import matplotlib
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from routers.analytics import compute_by_dimension_rows, compute_summary_values
from services.analytics_engine import filter_by_date_range
from services.analytics_repository import get_dataframe
from services.samsung_partner_config import SAMSUNG_PARTNER_LABELS, SAMSUNG_PARTNER_SOURCES

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from matplotlib.lines import Line2D  # noqa: E402
from matplotlib.patches import Patch  # noqa: E402
from matplotlib.ticker import FuncFormatter  # noqa: E402


PARTNER_DISPLAY = {
    **SAMSUNG_PARTNER_LABELS,
    "reliance": "Reliance ResQ",
    "godrej": "Godrej",
    "hitachi": "Hitachi",
}

PARTNER_LOGO = {
    "samsung_vs": "vs_logo.jpg",
    "samsung_croma": "croma_logo.jpg",
    "samsung_croma_dsdsg": "croma_logo.jpg",
    "samsung_reliance_digital": "reliance_digital_logo.png",
    "reliance": "resq.png",
    "godrej": "Group 1244833444.png",
    "hitachi": "hitachi_logo.png",
}

DEFAULT_PARTNERS = [*SAMSUNG_PARTNER_SOURCES, "reliance", "godrej", "hitachi"]

DIMENSION_CANDIDATES = {
    "samsung_vs": ["month", "state", "plan_category", "device_plan_category", "model_code"],
    "samsung_croma": ["month", "state", "plan_category", "device_plan_category", "model_code"],
    "samsung_croma_dsdsg": ["month", "state", "plan_category", "device_plan_category", "model_code"],
    "samsung_reliance_digital": ["month", "state", "plan_category", "device_plan_category", "model_code"],
    "reliance": ["month", "state", "plan_category", "article_brand"],
    "godrej": ["month", "state", "channel", "product_category"],
    "hitachi": ["month", "state", "channel", "product_category"],
}

WEEK_DATE_COLUMNS = ["Date", "Start_Date", "Start Date", "Plan Start Date", "Month"]
WEEK_VALUE_COLUMNS = {
    "gross_premium": ["Amount", "Gross Premium", "Plan Selling Price", "Customer Premium"],
}
VALID_WEEK_WINDOWS = {2, 3, 4, 6}
MODEL_DIMENSION_FALLBACKS = ["model_code", "model_code_1", "model", "product_model", "article_name"]
DEVICE_PLAN_CATEGORY_ORDER = [
    "Mass",
    "Mid",
    "High",
    "Premium",
    "Super Premium",
    "Luxury Flip",
    "Luxury Fold",
]
PLAN_CATEGORY_ORDER = ["SP", "COMBO", "ADLD", "EW"]

BACKGROUND_HEX = "EFEFEF"
TITLE_HEX = "1B2246"
TEXT_HEX = "263238"
ACCENT_HEX = "2D6BE8"
SECONDARY_HEX = "F97316"
PREVIEW_BG_HEX = "D9DCE2"
PANEL_BG_HEX = "D1D5DB"
PANEL_LINE_HEX = "C5CAD2"
CARD_BG_HEX = "F5F7FA"
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5
SLIDE_SAFE_MARGIN_IN = 0.18
TABLE_BOTTOM_SAFE_IN = 6.95


@dataclass
class DeckScope:
    source: str
    dataset_type: str
    job_id: str | None
    from_date: str | None
    to_date: str | None
    include_tables: bool
    week_window: int


def resolve_partners(raw_partners: list[str] | None) -> list[str]:
    if not raw_partners:
        return list(DEFAULT_PARTNERS)

    aliases = {
        "samsung": list(SAMSUNG_PARTNER_SOURCES),
        "samsung_vs": ["samsung_vs"],
        "samsung_vijay_sales": ["samsung_vs"],
        "vijay sales": ["samsung_vs"],
        "samsung_croma": ["samsung_croma"],
        "croma": ["samsung_croma"],
        "samsung protect max": ["samsung_croma"],
        "samsung protect max croma": ["samsung_croma"],
        "protect max": ["samsung_croma"],
        "protect max croma": ["samsung_croma"],
        "croma protect max": ["samsung_croma"],
        "samsung_croma_dsdsg": ["samsung_croma_dsdsg"],
        "croma ds dsg": ["samsung_croma_dsdsg"],
        "croma ds/dsg": ["samsung_croma_dsdsg"],
        "dsdsg": ["samsung_croma_dsdsg"],
        "ds dsg": ["samsung_croma_dsdsg"],
        "ds/dsg": ["samsung_croma_dsdsg"],
        "ds-dsg": ["samsung_croma_dsdsg"],
        "samsung_reliance_digital": ["samsung_reliance_digital"],
        "samsung reliance digital": ["samsung_reliance_digital"],
        "reliance digital": ["samsung_reliance_digital"],
        "reliance": ["reliance"],
        "reliance_resq": ["reliance"],
        "resq": ["reliance"],
        "godrej": ["godrej"],
        "goodrej": ["godrej"],
        "goddrej": ["godrej"],
        "hitachi": ["hitachi"],
    }

    ordered: list[str] = []
    seen: set[str] = set()
    for item in raw_partners:
        key = (item or "").strip().lower()
        expanded = aliases.get(key, [])
        for partner in expanded:
            if partner in PARTNER_DISPLAY and partner not in seen:
                seen.add(partner)
                ordered.append(partner)
    return ordered or list(DEFAULT_PARTNERS)


def generate_partner_deck_pptx(
    *,
    db: Session,
    partners: list[str],
    dataset_type: str,
    job_id: str | None,
    from_date: str | None,
    to_date: str | None,
    include_tables: bool,
    week_window: int,
) -> tuple[bytes, str]:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    normalized_dataset = "claims" if str(dataset_type).strip().lower() == "claims" else "sales"
    resolved_partners = resolve_partners(partners)
    normalized_week_window = _normalize_week_window(week_window)

    for source in resolved_partners:
        scope = DeckScope(
            source=source,
            dataset_type=normalized_dataset,
            job_id=(job_id or "").strip() or None,
            from_date=(from_date or "").strip() or None,
            to_date=(to_date or "").strip() or None,
            include_tables=include_tables,
            week_window=normalized_week_window,
        )
        _add_partner_section(prs, db=db, scope=scope)

    payload = BytesIO()
    prs.save(payload)
    payload.seek(0)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"partner_deck_{normalized_dataset}_{stamp}.pptx"
    return payload.getvalue(), filename


def build_partner_deck_preview(
    *,
    db: Session,
    partners: list[str],
    dataset_type: str,
    job_id: str | None,
    from_date: str | None,
    to_date: str | None,
    week_window: int,
) -> list[dict[str, Any]]:
    normalized_dataset = "claims" if str(dataset_type).strip().lower() == "claims" else "sales"
    resolved_partners = resolve_partners(partners)
    normalized_week_window = _normalize_week_window(week_window)

    preview_items: list[dict[str, Any]] = []
    for source in resolved_partners:
        scope = DeckScope(
            source=source,
            dataset_type=normalized_dataset,
            job_id=(job_id or "").strip() or None,
            from_date=(from_date or "").strip() or None,
            to_date=(to_date or "").strip() or None,
            include_tables=False,
            week_window=normalized_week_window,
        )
        preview_items.append(_build_partner_preview_item(db=db, scope=scope))
    return preview_items


def _build_partner_preview_item(*, db: Session, scope: DeckScope) -> dict[str, Any]:
    summary = compute_summary_values(
        db=db,
        job_id=scope.job_id,
        source=scope.source,
        dataset_type=scope.dataset_type,
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    summary_map = summary if isinstance(summary, dict) else {}

    trend_dimension = "month"
    gross_rows, gross_metric = _fetch_primary_metric_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="month",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    quantity_rows = _fetch_dimension_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="month",
        metric="quantity",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )

    if scope.dataset_type == "sales" and scope.source in SAMSUNG_PARTNER_SOURCES:
        week_gross = _fetch_samsung_week_rows(
            db=db,
            source=scope.source,
            dataset_type=scope.dataset_type,
            job_id=scope.job_id,
            metric="gross_premium",
            from_date=scope.from_date,
            to_date=scope.to_date,
        )
        week_qty = _fetch_samsung_week_rows(
            db=db,
            source=scope.source,
            dataset_type=scope.dataset_type,
            job_id=scope.job_id,
            metric="quantity",
            from_date=scope.from_date,
            to_date=scope.to_date,
        )
        week_gross = _limit_recent_weeks(week_gross, scope.week_window)
        week_qty = _limit_recent_weeks(week_qty, scope.week_window)
        if week_gross:
            trend_dimension = "week"
            gross_rows = week_gross
            gross_metric = "gross_premium"
            quantity_rows = week_qty

    gross_points = _to_primary_points(gross_rows, trend_dimension, gross_metric)
    quantity_points = _to_points(quantity_rows, trend_dimension, "quantity")

    state_gross_rows, state_metric = _fetch_primary_metric_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="state",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    state_qty_rows = _fetch_dimension_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="state",
        metric="quantity",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )

    state_points = _merge_metric_points(
        _to_primary_points(state_gross_rows, "state", state_metric),
        _to_points(state_qty_rows, "state", "quantity"),
        dimension="state",
    )
    state_points = sorted(state_points, key=lambda item: item.get("gross_premium", 0.0), reverse=True)[:8]

    product_points: list[dict[str, Any]] = []
    if scope.source in SAMSUNG_PARTNER_SOURCES:
        product_gross_rows, product_metric = _fetch_primary_metric_rows(
            db=db,
            source=scope.source,
            dataset_type=scope.dataset_type,
            job_id=scope.job_id,
            dimension="model_code",
            from_date=scope.from_date,
            to_date=scope.to_date,
        )
        product_qty_rows = _fetch_dimension_rows(
            db=db,
            source=scope.source,
            dataset_type=scope.dataset_type,
            job_id=scope.job_id,
            dimension="model_code",
            metric="quantity",
            from_date=scope.from_date,
            to_date=scope.to_date,
        )
        product_points = _merge_metric_points(
            _to_primary_points(product_gross_rows, "model_code", product_metric),
            _to_points(product_qty_rows, "model_code", "quantity"),
            dimension="model_code",
        )
        product_points = sorted(product_points, key=lambda item: item.get("gross_premium", 0.0), reverse=True)[:8]

    trend_points = _merge_metric_points(gross_points, quantity_points, dimension=trend_dimension)
    insights = _build_trend_insights(gross_rows, trend_dimension, gross_metric)

    return {
        "source": scope.source,
        "display_name": PARTNER_DISPLAY.get(scope.source, scope.source.title()),
        "logo": PARTNER_LOGO.get(scope.source),
        "summary": {
            "gross_premium": _to_float(summary_map.get("gross_premium")),
            "quantity": _to_float(summary_map.get("units_sold")),
        },
        "trend_dimension": trend_dimension,
        "trend_points": trend_points,
        "state_points": state_points,
        "product_points": product_points,
        "insights": insights,
    }


def _add_partner_section(prs: Presentation, *, db: Session, scope: DeckScope) -> None:
    preview = _build_partner_preview_item(db=db, scope=scope)
    summary_map = preview.get("summary", {}) if isinstance(preview.get("summary"), dict) else {}
    summary_payload = {
        "gross_premium": _to_float(summary_map.get("gross_premium")),
        "units_sold": _to_float(summary_map.get("quantity")),
    }
    month_rows = _build_month_trend_points(db=db, scope=scope)
    _add_summary_slide(prs, scope=scope, summary=summary_payload, month_rows=month_rows)

    partner_name = PARTNER_DISPLAY.get(scope.source, scope.source.title())

    if month_rows:
        _add_analysis_slide(
            prs,
            scope=scope,
            rows=_normalize_preview_rows(month_rows, max_rows=12),
            dimension="month",
            title=f"{partner_name}: Month Analysis",
            table_title="Top Month values",
            chart_mode="split_line",
        )

    trend_dim = str(preview.get("trend_dimension") or "month").strip().lower()
    trend_rows = _normalize_preview_rows(preview.get("trend_points"), max_rows=10)
    if trend_dim == "week" and trend_rows:
        _add_analysis_slide(
            prs,
            scope=scope,
            rows=trend_rows,
            dimension="week",
            title=f"{partner_name}: Week Analysis",
            table_title="Top Week values",
            chart_mode="split_line",
        )

    dimension_payloads: list[tuple[str, list[dict[str, Any]]]] = []
    dimension_rows_map: dict[str, list[dict[str, Any]]] = {}
    for dimension in DIMENSION_CANDIDATES.get(scope.source, []):
        if dimension == "month":
            continue
        rows = _build_dimension_points(db=db, scope=scope, dimension=dimension)
        if not rows:
            continue
        dimension_payloads.append((dimension, rows))
        dimension_rows_map[dimension] = rows

    plan_rows = dimension_rows_map.get("plan_category", [])
    device_plan_rows = dimension_rows_map.get("device_plan_category", [])
    if scope.dataset_type == "sales" and (plan_rows or device_plan_rows):
        _add_asp_dual_category_slide(
            prs,
            scope=scope,
            plan_rows=plan_rows,
            device_rows=device_plan_rows,
            title=f"{partner_name}: ASP Analysis (Category View)",
        )

    model_rows = dimension_rows_map.get("model_code", [])
    if model_rows:
        _add_product_model_dual_slide(
            prs,
            scope=scope,
            rows=model_rows,
            title=f"{partner_name}: Product Model Drilldown",
        )

    state_full_rows = _build_state_points_full(db=db, scope=scope)
    if state_full_rows:
        top_states = state_full_rows[:5]
        bottom_states = sorted(state_full_rows, key=lambda item: item.get("gross_premium", 0.0))[:5]
        _add_analysis_slide(
            prs,
            scope=scope,
            rows=top_states,
            dimension="state",
            title=f"{partner_name}: Top 5 State Contributors",
            table_title="Top 5 State Contributors",
            chart_mode="split_line",
        )
        _add_analysis_slide(
            prs,
            scope=scope,
            rows=bottom_states,
            dimension="state",
            title=f"{partner_name}: Bottom 5 State Contributors",
            table_title="Bottom 5 State Contributors",
            chart_mode="split_line",
        )

    small_dimensions = [(dim, rows) for dim, rows in dimension_payloads if len(rows) <= 6 and dim != "state"]
    used_small_dimensions: set[str] = set()
    if small_dimensions:
        selected_small = small_dimensions[:2]
        _add_dual_small_charts_slide(slide_deck=prs, scope=scope, charts=selected_small)
        used_small_dimensions = {dim for dim, _ in selected_small}

    for dimension, rows in dimension_payloads:
        if dimension == "model_code":
            continue
        if dimension in used_small_dimensions and dimension != "device_plan_category":
            continue
        readable = _prettify_label(dimension)
        table_title = f"Top {readable} values"
        chart_mode = "split_line" if dimension in {"plan_category", "device_plan_category"} else "combo"
        _add_analysis_slide(
            prs,
            scope=scope,
            rows=rows,
            dimension=dimension,
            title=f"{partner_name}: {readable} Analysis",
            table_title=table_title,
            chart_mode=chart_mode,
        )


def _build_month_trend_points(*, db: Session, scope: DeckScope) -> list[dict[str, Any]]:
    gross_rows, gross_metric = _fetch_primary_metric_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="month",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    qty_rows = _fetch_dimension_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="month",
        metric="quantity",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )

    month_points = _merge_metric_points(
        _to_primary_points(gross_rows, "month", gross_metric),
        _to_points(qty_rows, "month", "quantity"),
        dimension="month",
    )
    if len(month_points) > 6:
        month_points = month_points[-6:]
    return month_points


def _build_dimension_points(*, db: Session, scope: DeckScope, dimension: str) -> list[dict[str, Any]]:
    gross_rows, gross_metric = _fetch_primary_metric_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension=dimension,
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    qty_rows = _fetch_dimension_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension=dimension,
        metric="quantity",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )

    points = _merge_metric_points(
        _to_primary_points(gross_rows, dimension, gross_metric),
        _to_points(qty_rows, dimension, "quantity"),
        dimension=dimension,
    )
    if dimension == "device_plan_category":
        points = sorted(
            points,
            key=lambda item: (
                _device_plan_category_rank(str(item.get("label") or "")),
                -_to_float(item.get("gross_premium")),
            ),
        )
    elif dimension not in {"month", "week"}:
        points = sorted(points, key=lambda item: item.get("gross_premium", 0.0), reverse=True)[:10]
    return _normalize_preview_rows(points, max_rows=10)


def _build_state_points_full(*, db: Session, scope: DeckScope) -> list[dict[str, Any]]:
    gross_rows, gross_metric = _fetch_primary_metric_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="state",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )
    qty_rows = _fetch_dimension_rows(
        db=db,
        source=scope.source,
        dataset_type=scope.dataset_type,
        job_id=scope.job_id,
        dimension="state",
        metric="quantity",
        from_date=scope.from_date,
        to_date=scope.to_date,
    )

    points = _merge_metric_points(
        _to_primary_points(gross_rows, "state", gross_metric),
        _to_points(qty_rows, "state", "quantity"),
        dimension="state",
    )
    points = sorted(points, key=lambda item: item.get("gross_premium", 0.0), reverse=True)
    return _normalize_preview_rows(points, max_rows=20)


def _add_analysis_slide(
    slide_deck: Presentation,
    *,
    scope: DeckScope,
    rows: list[dict[str, Any]],
    dimension: str,
    title: str,
    table_title: str,
    chart_mode: str = "combo",
) -> None:
    if not rows:
        return

    slide = slide_deck.slides.add_slide(slide_deck.slide_layouts[6])
    _style_slide_background(slide)
    _add_partner_branding(slide, scope.source)
    _add_title_block(slide, title, _build_subtitle(scope))

    chart_rows = rows[:10]
    if chart_mode == "split_line":
        chart_img = _render_mom_vertical_split_chart_image(
            chart_rows,
            title=f"{_prettify_label(dimension)} Split",
            primary_label=_primary_metric_label(scope.dataset_type),
        )
    else:
        chart_img = _render_combo_chart_image(
            chart_rows,
            primary_label=_primary_metric_label(scope.dataset_type),
        )
    slide.shapes.add_picture(chart_img, Inches(0.65), Inches(1.6), Inches(8.0), Inches(4.35))

    metric_rows = [{"label": item.get("label"), "gross_premium": item.get("gross_premium")} for item in chart_rows]
    insights = _build_trend_insights(
        metric_rows,
        dimension,
        "gross_premium",
        metric_label=_primary_metric_label(scope.dataset_type),
    )
    _add_insight_box(slide, insights, left=Inches(8.8), top=Inches(1.6), width=Inches(3.95), height=Inches(4.35))

    if scope.include_tables:
        _add_data_table(
            slide,
            points=[(str(item.get("label") or ""), _to_float(item.get("gross_premium"))) for item in chart_rows],
            title=table_title,
            top_in=6.06,
            max_height_in=0.94,
            max_rows=3,
        )


def _build_asp_points(rows: list[dict[str, Any]], *, dimension: str) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    for row in rows:
        label = _canonical_dimension_label(str(row.get("label") or "").strip(), dimension)
        if not label:
            continue
        gross = _to_float(row.get("gross_premium"))
        qty = _to_float(row.get("quantity"))
        if qty <= 0:
            continue
        output.append({"label": label, "asp": gross / qty, "gross_premium": gross, "quantity": qty})
    if dimension == "plan_category":
        order_map = {label: idx for idx, label in enumerate(PLAN_CATEGORY_ORDER)}
        output = sorted(output, key=lambda row: (order_map.get(str(row.get("label") or ""), 99), -_to_float(row.get("asp"))))
    elif dimension == "device_plan_category":
        order_map = {label: idx for idx, label in enumerate(DEVICE_PLAN_CATEGORY_ORDER)}
        output = sorted(output, key=lambda row: (order_map.get(str(row.get("label") or ""), 99), -_to_float(row.get("asp"))))
    if len(output) > 10:
        output = output[:10]
    return output


def _add_asp_dual_category_slide(
    slide_deck: Presentation,
    *,
    scope: DeckScope,
    plan_rows: list[dict[str, Any]],
    device_rows: list[dict[str, Any]],
    title: str,
) -> None:
    plan_asp_rows = _build_asp_points(plan_rows, dimension="plan_category")[:8]
    device_asp_rows = _build_asp_points(device_rows, dimension="device_plan_category")[:8]
    if not plan_asp_rows and not device_asp_rows:
        return

    slide = slide_deck.slides.add_slide(slide_deck.slide_layouts[6])
    _style_slide_background(slide)
    _add_partner_branding(slide, scope.source)
    _add_title_block(slide, title, _build_subtitle(scope))

    left_rows = plan_asp_rows or device_asp_rows
    right_rows = device_asp_rows or plan_asp_rows

    left_img = _render_asp_chart_image(left_rows, dimension="plan_category", title="ASP by Plan Category")
    slide.shapes.add_picture(left_img, Inches(0.65), Inches(1.6), Inches(5.95), Inches(3.35))
    right_img = _render_asp_chart_image(right_rows, dimension="device_plan_category", title="ASP by Device Plan Category")
    slide.shapes.add_picture(right_img, Inches(6.8), Inches(1.6), Inches(5.95), Inches(3.35))

    plan_insights = _build_asp_insights(left_rows, "plan_category")
    device_insights = _build_asp_insights(right_rows, "device_plan_category")
    combined_insights = [
        f"Plan category ASP: {plan_insights[0]}" if plan_insights else "",
        f"Device plan ASP: {device_insights[0]}" if device_insights else "",
        f"Top ASP plan category: {plan_insights[2]}" if len(plan_insights) > 2 else "",
        f"Top ASP device plan category: {device_insights[2]}" if len(device_insights) > 2 else "",
    ]
    combined_insights = [line for line in combined_insights if line]
    _add_insight_box(slide, combined_insights, left=Inches(0.65), top=Inches(5.0), width=Inches(12.1), height=Inches(0.95))

    if scope.include_tables:
        _add_data_table(
            slide,
            points=[(str(item.get("label") or ""), _to_float(item.get("asp"))) for item in left_rows],
            title="ASP by Plan Category",
            top_in=6.06,
            max_height_in=0.90,
            max_rows=2,
        )


def _add_product_model_dual_slide(
    slide_deck: Presentation,
    *,
    scope: DeckScope,
    rows: list[dict[str, Any]],
    title: str,
) -> None:
    if not rows:
        return

    slide = slide_deck.slides.add_slide(slide_deck.slide_layouts[6])
    _style_slide_background(slide)
    _add_partner_branding(slide, scope.source)
    _add_title_block(slide, title, _build_subtitle(scope))

    chart_rows = rows[:10]
    gp_img = _render_metric_bar_chart_image(
        chart_rows,
        metric_key="gross_premium",
        title=f"{_primary_metric_label(scope.dataset_type)} by Product Model",
        color_hex=ACCENT_HEX,
        quantity_mode=False,
        primary_label=_primary_metric_label(scope.dataset_type),
    )
    qty_img = _render_metric_bar_chart_image(
        chart_rows,
        metric_key="quantity",
        title="No. of Plans by Product Model",
        color_hex=SECONDARY_HEX,
        quantity_mode=True,
    )

    slide.shapes.add_picture(gp_img, Inches(0.65), Inches(1.6), Inches(5.95), Inches(3.35))
    slide.shapes.add_picture(qty_img, Inches(6.8), Inches(1.6), Inches(5.95), Inches(3.35))

    metric_rows = [{"label": item.get("label"), "gross_premium": item.get("gross_premium")} for item in chart_rows]
    insights = _build_trend_insights(
        metric_rows,
        "model_code",
        "gross_premium",
        metric_label=_primary_metric_label(scope.dataset_type),
    )
    _add_insight_box(slide, insights, left=Inches(0.65), top=Inches(5.0), width=Inches(12.1), height=Inches(0.95))

    if scope.include_tables:
        _add_data_table(
            slide,
            points=[(str(item.get("label") or ""), _to_float(item.get("gross_premium"))) for item in chart_rows[:5]],
            title="Top Product Model values",
            top_in=6.06,
            max_height_in=0.90,
            max_rows=2,
        )


def _add_dual_small_charts_slide(
    *,
    slide_deck: Presentation,
    scope: DeckScope,
    charts: list[tuple[str, list[dict[str, Any]]]],
) -> None:
    if not charts:
        return

    slide = slide_deck.slides.add_slide(slide_deck.slide_layouts[6])
    _style_slide_background(slide)
    _add_partner_branding(slide, scope.source)

    partner_name = PARTNER_DISPLAY.get(scope.source, scope.source.title())
    _add_title_block(slide, f"{partner_name}: Multi-Graph Drilldown", _build_subtitle(scope))

    left_dimension, left_rows = charts[0]
    left_title = f"{_prettify_label(left_dimension)} (Split + Line)"
    left_img = _render_mom_vertical_split_chart_image(
        left_rows,
        title=left_title,
        primary_label=_primary_metric_label(scope.dataset_type),
    )
    slide.shapes.add_picture(left_img, Inches(0.65), Inches(1.55), Inches(5.85), Inches(4.25))

    right_rows = left_rows
    right_dimension = left_dimension
    if len(charts) > 1:
        right_dimension, right_rows = charts[1]
    right_title = f"{_prettify_label(right_dimension)} (Pie)"
    right_img = _render_pie_chart_image(right_rows, right_title)
    slide.shapes.add_picture(right_img, Inches(6.85), Inches(1.55), Inches(5.85), Inches(4.25))

    left_insights = _build_trend_insights(
        [{"label": row.get("label"), "gross_premium": row.get("gross_premium")} for row in left_rows],
        left_dimension,
        "gross_premium",
        metric_label=_primary_metric_label(scope.dataset_type),
    )
    right_insights = _build_trend_insights(
        [{"label": row.get("label"), "gross_premium": row.get("gross_premium")} for row in right_rows],
        right_dimension,
        "gross_premium",
        metric_label=_primary_metric_label(scope.dataset_type),
    )
    merged_insights = [
        f"{_prettify_label(left_dimension)}: {left_insights[0]}" if left_insights else "",
        f"{_prettify_label(right_dimension)}: {right_insights[0]}" if right_insights else "",
        f"{_primary_metric_label(scope.dataset_type)} and Quantity only. Earned premium metrics are excluded from the deck.",
    ]
    merged_insights = [line for line in merged_insights if line]
    _add_insight_box(slide, merged_insights, left=Inches(0.65), top=Inches(5.95), width=Inches(12.1), height=Inches(1.2))


def _add_partner_cover_slide(prs: Presentation, *, scope: DeckScope, preview: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_preview_background(slide)
    _add_partner_branding(slide, scope.source)

    display_name = str(preview.get("display_name") or PARTNER_DISPLAY.get(scope.source, scope.source.title()))
    subtitle = _build_subtitle(scope)

    _add_rect_card(slide, left=0.62, top=1.12, width=12.1, height=4.95, fill_hex=PANEL_BG_HEX, line_hex=PANEL_LINE_HEX)

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(10.8), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{display_name} Partner Deck"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(TITLE_HEX)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.62), Inches(10.8), Inches(0.45))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor.from_string("5F6C85")
    p2.alignment = PP_ALIGN.CENTER

    metric_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.35), Inches(10.0), Inches(1.2))
    mt = metric_box.text_frame
    mt.clear()
    mp = mt.paragraphs[0]
    mp.text = f"Deck Metrics: {_primary_metric_label(scope.dataset_type)} + Quantity"
    mp.font.size = Pt(23)
    mp.font.bold = True
    mp.font.color.rgb = RGBColor.from_string(ACCENT_HEX)
    mp.alignment = PP_ALIGN.CENTER

    ip = mt.add_paragraph()
    ip.text = "Preview-aligned slide format with trend, drilldown, and AI insight sections."
    ip.font.size = Pt(13)
    ip.font.color.rgb = RGBColor.from_string(TEXT_HEX)
    ip.alignment = PP_ALIGN.CENTER
    ip.space_before = Pt(8)


def _add_partner_preview_main_slide(prs: Presentation, *, scope: DeckScope, preview: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_preview_background(slide)
    _add_preview_header(slide, scope=scope, preview=preview)

    trend_rows = _normalize_preview_rows(preview.get("trend_points"), max_rows=12)
    state_rows = _normalize_preview_rows(preview.get("state_points"), max_rows=8)
    insights = [str(item) for item in preview.get("insights", []) if str(item).strip()]

    # Panel shells
    _add_rect_card(slide, left=0.45, top=0.92, width=3.95, height=3.35, fill_hex=PANEL_BG_HEX, line_hex=PANEL_LINE_HEX)
    _add_rect_card(slide, left=4.52, top=0.92, width=8.35, height=3.35, fill_hex=PANEL_BG_HEX, line_hex=PANEL_LINE_HEX)
    _add_rect_card(slide, left=0.45, top=4.4, width=12.42, height=2.75, fill_hex=PANEL_BG_HEX, line_hex=PANEL_LINE_HEX)

    # Slide 1 snapshot card
    _add_small_heading(slide, left=0.67, top=1.08, width=3.4, text="SLIDE 1: SNAPSHOT")
    _add_rect_card(slide, left=0.66, top=1.46, width=1.58, height=0.86, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")
    _add_rect_card(slide, left=2.31, top=1.46, width=1.58, height=0.86, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")

    summary = preview.get("summary", {}) if isinstance(preview.get("summary"), dict) else {}
    gross = _to_float(summary.get("gross_premium"))
    qty = _to_float(summary.get("quantity"))
    _add_metric_card_text(
        slide,
        left=0.79,
        top=1.58,
        title=_primary_metric_label(scope.dataset_type).upper(),
        value=f"Rs {_compact_number(gross)}",
    )
    _add_metric_card_text(slide, left=2.44, top=1.58, title="QUANTITY", value=_format_quantity_number(qty))

    _add_rect_card(slide, left=0.66, top=2.43, width=3.23, height=1.69, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")
    _add_small_heading(slide, left=0.8, top=2.56, width=3.0, text="INSIGHTS")
    _add_insight_list_block(slide, left=0.8, top=2.84, width=2.95, height=1.2, insights=insights[:5], font_size=8.9)

    # Slide 2 trend chart
    trend_dim = str(preview.get("trend_dimension") or "month").strip().lower()
    trend_label = "WEEK-WISE" if trend_dim == "week" else "MONTH-WISE"
    _add_small_heading(slide, left=4.74, top=1.08, width=7.8, text=f"SLIDE 2: TREND ({trend_label})")
    _add_rect_card(slide, left=4.73, top=1.46, width=7.95, height=2.6, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")

    trend_img = _render_combo_chart_image(
        trend_rows,
        primary_label=_primary_metric_label(scope.dataset_type),
    )
    slide.shapes.add_picture(trend_img, Inches(5.02), Inches(1.66), Inches(7.34), Inches(2.15))

    # Slide 3 state drilldown
    _add_small_heading(slide, left=0.67, top=4.58, width=4.2, text="SLIDE 3: STATE DRILLDOWN")
    _add_rect_card(slide, left=0.66, top=4.95, width=12.0, height=2.0, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")

    state_img = _render_combo_chart_image(
        state_rows,
        primary_label=_primary_metric_label(scope.dataset_type),
    )
    slide.shapes.add_picture(state_img, Inches(0.95), Inches(5.08), Inches(11.45), Inches(1.73))

    if scope.include_tables and state_rows:
        _add_compact_table_summary(
            slide,
            rows=state_rows[:4],
            left=Inches(0.9),
            top=Inches(6.84),
            width=Inches(11.5),
            title="State Table",
        )


def _add_partner_product_slide(prs: Presentation, *, scope: DeckScope, preview: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _style_preview_background(slide)
    _add_preview_header(slide, scope=scope, preview=preview)

    product_rows = _normalize_preview_rows(preview.get("product_points"), max_rows=10)
    if not product_rows:
        return

    insights = [str(item) for item in preview.get("insights", []) if str(item).strip()]

    _add_rect_card(slide, left=0.45, top=0.92, width=12.42, height=6.18, fill_hex=PANEL_BG_HEX, line_hex=PANEL_LINE_HEX)
    _add_small_heading(slide, left=0.7, top=1.08, width=6.4, text="SLIDE 4: PRODUCT DRILLDOWN (A17 / FOLD 6 / ETC.)")

    _add_rect_card(slide, left=0.69, top=1.45, width=3.6, height=5.3, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")
    _add_small_heading(slide, left=0.86, top=1.62, width=3.1, text="INSIGHTS")
    _add_insight_list_block(slide, left=0.86, top=1.9, width=3.1, height=4.7, insights=insights[:7], font_size=9.3)

    _add_rect_card(slide, left=4.45, top=1.45, width=8.15, height=5.3, fill_hex=CARD_BG_HEX, line_hex="E3E7EC")
    product_img = _render_combo_chart_image(
        product_rows,
        primary_label=_primary_metric_label(scope.dataset_type),
    )
    slide.shapes.add_picture(product_img, Inches(4.75), Inches(1.72), Inches(7.55), Inches(4.75))

    if scope.include_tables:
        _add_compact_table_summary(
            slide,
            rows=product_rows[:5],
            left=Inches(4.75),
            top=Inches(6.52),
            width=Inches(7.55),
            title="Product Table",
        )


def _add_preview_header(slide, *, scope: DeckScope, preview: dict[str, Any]) -> None:
    _add_partner_branding(slide, scope.source)
    display_name = str(preview.get("display_name") or PARTNER_DISPLAY.get(scope.source, scope.source.title()))
    trend_dim = str(preview.get("trend_dimension") or "month").lower().strip()
    week_note = f" (last {scope.week_window} weeks)" if trend_dim == "week" else ""

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(8.8), Inches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = display_name
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(TITLE_HEX)

    sub_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.57), Inches(10.2), Inches(0.28))
    st = sub_box.text_frame
    st.clear()
    sp = st.paragraphs[0]
    sp.text = f"Live slide view{week_note}. Changes here are reflected in download output."
    sp.font.size = Pt(11.5)
    sp.font.color.rgb = RGBColor.from_string("4F6483")


def _style_preview_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(PREVIEW_BG_HEX)


def _add_rect_card(slide, *, left: float, top: float, width: float, height: float, fill_hex: str, line_hex: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    shape.line.color.rgb = RGBColor.from_string(line_hex)
    shape.line.width = Pt(1)


def _add_small_heading(slide, *, left: float, top: float, width: float, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string("6A7B97")


def _add_metric_card_text(slide, *, left: float, top: float, title: str, value: str) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(1.3), Inches(0.66))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(9.2)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string("627999")

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor.from_string("1A2A49")
    p2.space_before = Pt(4)


def _add_insight_list_block(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    insights: list[str],
    font_size: float,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, insight in enumerate(insights):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = insight
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string("4A586F")
        p.space_after = Pt(4)


def _normalize_preview_rows(value: Any, *, max_rows: int) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    rows: list[dict[str, Any]] = []
    for row in value:
        if not isinstance(row, dict):
            continue
        label = str(row.get("label") or "").strip()
        if not label:
            continue
        rows.append(
            {
                "label": _short_label(label, max_len=18),
                "gross_premium": _to_float(row.get("gross_premium")),
                "claims": _to_float(row.get("claims")) or _to_float(row.get("gross_premium")),
                "quantity": _to_float(row.get("quantity")),
            }
        )
    return rows[:max_rows]


def _add_compact_table_summary(
    slide,
    *,
    rows: list[dict[str, Any]],
    left,
    top,
    width,
    title: str,
) -> None:
    if not rows:
        return
    left_in = float(getattr(left, "inches", 0.0))
    top_in = float(getattr(top, "inches", 0.0))
    width_in = float(getattr(width, "inches", 0.0))
    row_height_in = 0.11
    title_height_in = 0.14
    gap_in = 0.02
    min_table_height_in = row_height_in * 2  # header + at least one row
    min_total_height_in = title_height_in + gap_in + min_table_height_in

    safe_left, safe_top, safe_width, _ = _clamp_box_to_slide(
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_in=min_total_height_in,
    )

    title_box = slide.shapes.add_textbox(Inches(safe_left), Inches(safe_top), Inches(safe_width), Inches(title_height_in))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(8.6)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor.from_string("4A586F")
    max_bottom_in = min(TABLE_BOTTOM_SAFE_IN, SLIDE_HEIGHT_IN - SLIDE_SAFE_MARGIN_IN)
    min_total_height_in = title_height_in + gap_in + (row_height_in * 2)
    max_top_in = max(SLIDE_SAFE_MARGIN_IN, max_bottom_in - min_total_height_in)
    if safe_top > max_top_in:
        safe_top = max_top_in

    table_top_in = safe_top + title_height_in + gap_in
    available_height_in = max(0.22, max_bottom_in - table_top_in)

    clipped = rows[:5]
    row_count = len(clipped) + 1
    while row_count > 2 and (row_count * row_height_in) > available_height_in:
        clipped = clipped[:-1]
        row_count = len(clipped) + 1
    if not clipped:
        return

    table_shape = slide.shapes.add_table(
        row_count,
        2,
        Inches(safe_left),
        Inches(table_top_in),
        Inches(safe_width),
        Inches(min(row_count * row_height_in, available_height_in)),
    )
    table = table_shape.table
    col1 = safe_width * 0.68
    col2 = max(safe_width - col1, 1.2)
    table.columns[0].width = Inches(col1)
    table.columns[1].width = Inches(col2)

    table.cell(0, 0).text = "Dimension"
    table.cell(0, 1).text = "Value"
    _style_table_cell(table.cell(0, 0), bold=True)
    _style_table_cell(table.cell(0, 1), bold=True)

    for idx, row in enumerate(clipped, start=1):
        label = str(row.get("label") or "")
        value = _to_float(row.get("gross_premium"))
        table.cell(idx, 0).text = _short_label(label, max_len=20)
        table.cell(idx, 1).text = _compact_number(value)
        _style_table_cell(table.cell(idx, 0))
        _style_table_cell(table.cell(idx, 1), align_right=True)

    for idx in range(row_count):
        table.rows[idx].height = Inches(row_height_in)


def _fetch_dimension_rows(
    *,
    db: Session,
    source: str,
    dataset_type: str,
    job_id: str | None,
    dimension: str,
    metric: str,
    from_date: str | None,
    to_date: str | None,
) -> list[dict[str, Any]]:
    dimensions_to_try = [dimension]
    if dimension == "model_code":
        dimensions_to_try = MODEL_DIMENSION_FALLBACKS

    for candidate_dimension in dimensions_to_try:
        try:
            rows = compute_by_dimension_rows(
                db=db,
                job_id=job_id,
                dimension=candidate_dimension,
                metric=metric,
                source=source,
                dataset_type=dataset_type,
                bucket="month" if candidate_dimension == "month" else None,
                from_date=from_date,
                to_date=to_date,
                category_filters=None,
            )
        except Exception:
            continue

        normalized = [r for r in rows or [] if isinstance(r, dict)]
        if normalized:
            return normalized
    return []


def _primary_metric_key(dataset_type: str) -> str:
    return "claims" if str(dataset_type).strip().lower() == "claims" else "gross_premium"


def _primary_metric_label(dataset_type: str) -> str:
    return "Claims" if str(dataset_type).strip().lower() == "claims" else "Gross Premium"


def _fetch_primary_metric_rows(
    *,
    db: Session,
    source: str,
    dataset_type: str,
    job_id: str | None,
    dimension: str,
    from_date: str | None,
    to_date: str | None,
) -> tuple[list[dict[str, Any]], str]:
    primary_metric = _primary_metric_key(dataset_type)
    primary_rows = _fetch_dimension_rows(
        db=db,
        source=source,
        dataset_type=dataset_type,
        job_id=job_id,
        dimension=dimension,
        metric=primary_metric,
        from_date=from_date,
        to_date=to_date,
    )
    if primary_rows:
        return primary_rows, primary_metric

    if primary_metric != "gross_premium":
        fallback_rows = _fetch_dimension_rows(
            db=db,
            source=source,
            dataset_type=dataset_type,
            job_id=job_id,
            dimension=dimension,
            metric="gross_premium",
            from_date=from_date,
            to_date=to_date,
        )
        if fallback_rows:
            return fallback_rows, "gross_premium"

    return [], primary_metric


def _to_primary_points(
    rows: list[dict[str, Any]],
    dimension: str,
    metric_key: str,
) -> list[tuple[str, float]]:
    # Keep merge inputs in tuple form to avoid dict-unpack runtime failures.
    return _to_points(rows, dimension, metric_key)


def _normalize_week_window(value: int | None) -> int:
    try:
        parsed = int(value or 4)
    except Exception:
        parsed = 4
    return parsed if parsed in VALID_WEEK_WINDOWS else 4


def _limit_recent_weeks(rows: list[dict[str, Any]], week_window: int) -> list[dict[str, Any]]:
    if not rows:
        return []
    limit = _normalize_week_window(week_window)
    sorted_rows = sorted(rows, key=lambda item: _week_sort_key(str(item.get("week") or "")))
    return sorted_rows[-limit:]


def _week_sort_key(label: str) -> tuple[int, int, int]:
    text = str(label or "").strip()
    if not text:
        return (9999, 12, 31)

    first = text.split("-")[0].strip()
    ts = pd.to_datetime(first, errors="coerce")
    if pd.notna(ts):
        return (int(ts.year), int(ts.month), int(ts.day))

    month_map = {
        "jan": 1,
        "feb": 2,
        "mar": 3,
        "apr": 4,
        "may": 5,
        "jun": 6,
        "jul": 7,
        "aug": 8,
        "sep": 9,
        "oct": 10,
        "nov": 11,
        "dec": 12,
    }
    m = re.match(r"^(?P<day>\d{1,2})\s+(?P<mon>[A-Za-z]{3,9})$", first)
    if m:
        day = int(m.group("day"))
        mon = month_map.get(m.group("mon").lower()[:3], 12)
        return (2099, mon, day)

    return (9999, 12, 31)


def _fetch_samsung_week_rows(
    *,
    db: Session,
    source: str,
    dataset_type: str,
    job_id: str | None,
    metric: str,
    from_date: str | None,
    to_date: str | None,
) -> list[dict[str, Any]]:
    if dataset_type != "sales" or source not in SAMSUNG_PARTNER_SOURCES:
        return []

    try:
        df = get_dataframe(
            db=db,
            job_id=job_id,
            source=source,
            dataset_type="sales",
        )
    except Exception:
        return []

    if df is None or df.empty:
        return []

    df = filter_by_date_range(df, "sales", from_date, to_date)
    if df is None or df.empty:
        return []

    frame = df.copy()
    coalesced_week_date = pd.Series(pd.NaT, index=frame.index, dtype="datetime64[ns]")
    for candidate in WEEK_DATE_COLUMNS:
        col = candidate if candidate in frame.columns else _pick_first_column(frame, [candidate])
        if col is None or col not in frame.columns:
            continue
        parsed = _parse_date_like_series(frame[col])
        if parsed.isna().all():
            continue
        coalesced_week_date = coalesced_week_date.where(coalesced_week_date.notna(), parsed)
        if coalesced_week_date.notna().all():
            break

    frame["_week_date"] = coalesced_week_date
    frame = frame[frame["_week_date"].notna()].copy()
    if frame.empty:
        return []

    if _safe_key(metric) == "quantity":
        frame["_metric"] = 1.0
    else:
        frame["_metric"] = _resolve_metric_series(frame, metric)
    frame["_metric"] = pd.to_numeric(frame["_metric"], errors="coerce").fillna(0.0)

    week_start = frame["_week_date"] - pd.to_timedelta(frame["_week_date"].dt.weekday, unit="D")
    frame["_week_start"] = week_start.dt.normalize()

    grouped = (
        frame.groupby("_week_start", dropna=False)["_metric"]
        .sum()
        .reset_index()
        .sort_values("_week_start")
    )
    if grouped.empty:
        return []

    out: list[dict[str, Any]] = []
    for _, row in grouped.iterrows():
        start = pd.to_datetime(row["_week_start"], errors="coerce")
        if pd.isna(start):
            continue
        label = _format_week_label(start)
        out.append({"week": label, metric: _to_float(row["_metric"])})
    return out


def _pick_first_column(df: pd.DataFrame, candidates: list[str]) -> str | None:
    if df is None or df.empty:
        return None

    for candidate in candidates:
        if candidate in df.columns:
            return candidate

    target = [_safe_key(c) for c in candidates]
    for col in df.columns:
        if _safe_key(str(col)) in target:
            return str(col)
    return None


def _parse_date_like_series(series: pd.Series) -> pd.Series:
    cleaned = series.astype(str).str.strip().replace({"": pd.NA})
    direct = pd.to_datetime(cleaned, errors="coerce")
    if direct.notna().any():
        return direct

    m = cleaned.str.extract(r"^(?P<mon>[A-Za-z]{3,9})[-/ ](?P<yr>\d{2,4})$")
    if m.notna().any().any():
        month_map = {
            "jan": 1,
            "feb": 2,
            "mar": 3,
            "apr": 4,
            "may": 5,
            "jun": 6,
            "jul": 7,
            "aug": 8,
            "sep": 9,
            "oct": 10,
            "nov": 11,
            "dec": 12,
        }
        month_num = m["mon"].str.lower().str.slice(0, 3).map(month_map)
        year_num = pd.to_numeric(m["yr"], errors="coerce")
        is_two_digit = m["yr"].fillna("").str.len().eq(2)
        year_num = year_num.where(~is_two_digit, year_num + 2000)
        explicit = pd.to_datetime({"year": year_num, "month": month_num, "day": 1}, errors="coerce")
        direct = direct.fillna(explicit)

    return direct


def _resolve_metric_series(df: pd.DataFrame, metric: str) -> pd.Series:
    normalized = _safe_key(metric)
    if normalized == "gross_premium":
        col = _pick_first_column(df, WEEK_VALUE_COLUMNS["gross_premium"])
        return pd.to_numeric(df[col], errors="coerce").fillna(0.0) if col else pd.Series(0.0, index=df.index)
    if normalized == "quantity":
        return pd.Series(1.0, index=df.index)
    if normalized == "claims":
        col = _pick_first_column(df, ["Net Amount", "Net_Amount", "Claim_Amount", "Claims Cost", "Claim Cost"])
        return pd.to_numeric(df[col], errors="coerce").fillna(0.0) if col else pd.Series(0.0, index=df.index)
    return pd.Series(0.0, index=df.index)


def _format_week_label(week_start: pd.Timestamp) -> str:
    start = pd.Timestamp(week_start).normalize()
    end = start + pd.Timedelta(days=6)
    return f"{start.strftime('%d %b')}-{end.strftime('%d %b')}"


def _style_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(BACKGROUND_HEX)


def _add_title_block(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(9.3), Inches(0.5))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(TITLE_HEX)

    subtitle_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.83), Inches(9.3), Inches(0.32))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    p2 = subtitle_tf.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor.from_string(TEXT_HEX)

    line = slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0.55),
        top=Inches(1.15),
        width=Inches(12.2),
        height=Pt(1.2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string("C7C7C7")
    line.line.fill.background()


def _build_subtitle(scope: DeckScope) -> str:
    date_label = "Full available date range"
    if scope.from_date and scope.to_date:
        date_label = f"{scope.from_date} to {scope.to_date}"
    elif scope.from_date:
        date_label = f"From {scope.from_date}"
    elif scope.to_date:
        date_label = f"Until {scope.to_date}"
    return f"{_prettify_label(scope.dataset_type)} dataset | {date_label}"


def _add_partner_branding(slide, source: str) -> None:
    repo_root = Path(__file__).resolve().parents[2]
    public_dir = repo_root / "frontend" / "my-app" / "public"

    logo_paths = []
    partner_logo = PARTNER_LOGO.get(source)
    if partner_logo:
        logo_paths.append(public_dir / partner_logo)
    logo_paths.append(public_dir / "Zopper Logo Original 1.png")

    x = 10.65
    for path in logo_paths:
        if not path.exists():
            continue
        try:
            slide.shapes.add_picture(str(path), Inches(x), Inches(0.23), height=Inches(0.42))
            x += 1.35
        except Exception:
            continue


def _add_summary_slide(
    slide_deck: Presentation,
    *,
    scope: DeckScope,
    summary: dict[str, Any],
    month_rows: list[dict[str, Any]],
) -> None:
    slide = slide_deck.slides.add_slide(slide_deck.slide_layouts[6])
    _style_slide_background(slide)
    _add_partner_branding(slide, scope.source)

    title = f"{PARTNER_DISPLAY.get(scope.source, scope.source.title())} - Executive Snapshot"
    subtitle = _build_subtitle(scope)
    _add_title_block(slide, title, subtitle)

    primary_label = _primary_metric_label(scope.dataset_type)
    cards = [
        (primary_label, float(summary.get("gross_premium", 0) or 0)),
        ("Quantity", float(summary.get("units_sold", 0) or 0)),
    ]

    left = 0.7
    for label, value in cards:
        box = slide.shapes.add_textbox(Inches(left), Inches(1.55), Inches(4.4), Inches(1.15))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(TEXT_HEX)

        p2 = tf.add_paragraph()
        p2.text = _format_metric_value(label, value)
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor.from_string(TITLE_HEX)
        p2.space_before = Pt(8)

        left += 4.7

    chart_rows = _normalize_preview_rows(month_rows, max_rows=6)
    chart_img = _render_mom_vertical_split_chart_image(chart_rows, primary_label=primary_label)
    slide.shapes.add_picture(chart_img, Inches(0.65), Inches(2.95), Inches(6.3), Inches(3.65))

    insight_rows = [{"label": row.get("label"), "gross_premium": row.get("gross_premium")} for row in chart_rows]
    insights = _build_trend_insights(
        insight_rows,
        "month",
        "gross_premium",
        metric_label=primary_label,
    )
    _add_insight_box(slide, insights, left=Inches(7.15), top=Inches(2.95), width=Inches(5.5), height=Inches(3.65))


def _add_insight_box(
    slide,
    insights: list[str],
    *,
    left=Inches(0.65),
    top=Inches(4.65),
    width=Inches(12.1),
    height=Inches(1.7),
) -> None:
    height_in = float(getattr(height, "inches", 1.7))
    compact_mode = height_in < 1.2
    heading_size = 11.5 if compact_mode else 13
    body_size = 8.9 if compact_mode else 10.5
    max_items = 3 if compact_mode else 7

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
    panel.line.color.rgb = RGBColor.from_string("D9E1EC")
    panel.line.width = Pt(1)

    box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), height - Inches(0.16))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    head = tf.paragraphs[0]
    head.text = "Insights"
    head.font.size = Pt(heading_size)
    head.font.bold = True
    head.font.color.rgb = RGBColor.from_string(TITLE_HEX)
    head.space_after = Pt(5)

    for item in insights[:max_items]:
        p = tf.add_paragraph()
        p.text = f"- {item}"
        p.font.size = Pt(body_size)
        p.font.color.rgb = RGBColor.from_string(TEXT_HEX)
        p.level = 0
        p.space_after = Pt(2)


def _add_data_table(
    slide,
    *,
    points: list[tuple[str, float]],
    title: str,
    left_in: float = 0.65,
    top_in: float = 6.2,
    width_in: float = 12.1,
    max_height_in: float = 1.12,
    max_rows: int = 4,
) -> None:
    if not points:
        return
    safe_max_rows = max(1, int(max_rows or 1))
    title_height_in = 0.17
    gap_in = 0.04
    row_height_in = 0.13
    min_table_height_in = row_height_in * 2
    min_total_height_in = title_height_in + gap_in + min_table_height_in

    safe_left, safe_top, safe_width, safe_max_height = _clamp_box_to_slide(
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_in=max(max_height_in, min_total_height_in),
    )
    max_bottom_in = min(TABLE_BOTTOM_SAFE_IN, SLIDE_HEIGHT_IN - SLIDE_SAFE_MARGIN_IN)
    max_allowed_height = max(0.2, max_bottom_in - safe_top)
    if safe_max_height > max_allowed_height:
        safe_max_height = max_allowed_height
    if safe_max_height < min_total_height_in:
        safe_top = max(SLIDE_SAFE_MARGIN_IN, max_bottom_in - min_total_height_in)
        safe_max_height = max(0.2, max_bottom_in - safe_top)

    table_top_in = safe_top + title_height_in + gap_in
    available_table_height_in = max(0.14, safe_max_height - (title_height_in + gap_in))

    clipped_points = points[:safe_max_rows]
    if not clipped_points:
        return

    table_rows = len(clipped_points) + 1
    while table_rows > 2 and (table_rows * row_height_in) > available_table_height_in:
        clipped_points = clipped_points[:-1]
        table_rows = len(clipped_points) + 1
    if not clipped_points:
        clipped_points = points[:1]
        table_rows = 2

    table_height_in = min(table_rows * row_height_in, available_table_height_in)

    title_box = slide.shapes.add_textbox(Inches(safe_left), Inches(safe_top), Inches(safe_width), Inches(title_height_in))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(TITLE_HEX)

    table_shape = slide.shapes.add_table(
        table_rows,
        2,
        Inches(safe_left),
        Inches(table_top_in),
        Inches(safe_width),
        Inches(table_height_in),
    )
    table = table_shape.table
    first_col_width = safe_width * 0.72
    second_col_width = max(safe_width - first_col_width, 1.4)
    table.columns[0].width = Inches(first_col_width)
    table.columns[1].width = Inches(second_col_width)

    table.cell(0, 0).text = "Dimension"
    table.cell(0, 1).text = "Value"
    _style_table_cell(table.cell(0, 0), bold=True)
    _style_table_cell(table.cell(0, 1), bold=True)

    for idx, (label, value) in enumerate(clipped_points[: table_rows - 1], start=1):
        table.cell(idx, 0).text = _short_label(str(label), max_len=28)
        table.cell(idx, 1).text = _compact_number(value)
        _style_table_cell(table.cell(idx, 0))
        _style_table_cell(table.cell(idx, 1), align_right=True)

    for idx in range(table_rows):
        table.rows[idx].height = Inches(row_height_in)


def _style_table_cell(cell, *, bold: bool = False, align_right: bool = False) -> None:
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string("F6F6F6" if bold else "FFFFFF")
    tf = cell.text_frame
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(7.8)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor.from_string(TEXT_HEX)
        paragraph.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT
        paragraph.space_after = Pt(0)


def _clamp_box_to_slide(*, left_in: float, top_in: float, width_in: float, height_in: float) -> tuple[float, float, float, float]:
    safe_left = max(SLIDE_SAFE_MARGIN_IN, float(left_in or 0.0))
    safe_top = max(SLIDE_SAFE_MARGIN_IN, float(top_in or 0.0))
    safe_width = max(1.4, float(width_in or 0.0))
    safe_height = max(0.2, float(height_in or 0.0))

    max_width = max(1.4, SLIDE_WIDTH_IN - safe_left - SLIDE_SAFE_MARGIN_IN)
    if safe_width > max_width:
        safe_width = max_width

    max_height = max(0.2, SLIDE_HEIGHT_IN - safe_top - SLIDE_SAFE_MARGIN_IN)
    if safe_height > max_height:
        safe_height = max_height

    # If top is too low for requested minimum height, shift upward.
    min_required = min(safe_height, SLIDE_HEIGHT_IN - (2 * SLIDE_SAFE_MARGIN_IN))
    max_top = SLIDE_HEIGHT_IN - SLIDE_SAFE_MARGIN_IN - min_required
    if safe_top > max_top:
        safe_top = max(SLIDE_SAFE_MARGIN_IN, max_top)
        max_height = max(0.2, SLIDE_HEIGHT_IN - safe_top - SLIDE_SAFE_MARGIN_IN)
        safe_height = min(safe_height, max_height)

    return safe_left, safe_top, safe_width, safe_height


def _render_chart_image(
    *,
    points: list[tuple[str, float]],
    title: str,
    kind: str,
    color: str,
) -> BytesIO:
    labels = [str(p[0]) for p in points]
    values = [float(p[1]) for p in points]
    if not labels:
        labels = ["No data"]
        values = [0.0]

    fig, ax = plt.subplots(figsize=(6.5, 3.2), dpi=120)
    fig.patch.set_facecolor("#EFEFEF")
    ax.set_facecolor("#EFEFEF")

    if kind == "line":
        ax.plot(range(len(values)), values, color=color, marker="o", linewidth=2.4, markersize=4)
        ax.fill_between(range(len(values)), values, color=color, alpha=0.12)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels([_short_label(l) for l in labels], rotation=0, fontsize=8)
    elif kind == "pie":
        safe_values = [v for v in values if v > 0]
        safe_labels = [labels[i] for i, v in enumerate(values) if v > 0]
        if not safe_values:
            safe_values = [1]
            safe_labels = ["No data"]
        ax.pie(
            safe_values[:8],
            labels=[_short_label(v, max_len=14) for v in safe_labels[:8]],
            autopct="%1.0f%%",
            startangle=90,
            colors=plt.cm.Pastel1.colors,
            wedgeprops={"linewidth": 0.8, "edgecolor": "#EFEFEF"},
            textprops={"fontsize": 8},
        )
        ax.axis("equal")
    else:
        # Use horizontal bars for dense categorical drilldowns.
        order = sorted(zip(labels, values), key=lambda item: item[1], reverse=True)[:10]
        cat = [o[0] for o in order]
        val = [o[1] for o in order]
        ax.barh(range(len(val)), val, color=color, alpha=0.86)
        ax.set_yticks(range(len(cat)))
        ax.set_yticklabels([_short_label(c) for c in cat], fontsize=8)
        ax.invert_yaxis()

    ax.set_title(title, fontsize=11, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=8)
    ax.tick_params(axis="x", labelsize=8)
    ax.grid(axis="x", linestyle="-", linewidth=0.4, alpha=0.25)
    for side in ("top", "right", "left"):
        ax.spines[side].set_visible(False)
    ax.spines["bottom"].set_alpha(0.3)

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=120, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_mom_vertical_split_chart_image(
    rows: list[dict[str, Any]],
    *,
    title: str = "Month-on-Month Split",
    primary_label: str = "Gross Premium",
) -> BytesIO:
    labels = [str(row.get("label") or "") for row in rows]
    gross = [_to_float(row.get("gross_premium")) for row in rows]
    quantity = [_to_float(row.get("quantity")) for row in rows]

    if not labels:
        labels = ["No data"]
        gross = [0.0]
        quantity = [0.0]

    x = list(range(len(labels)))
    width = 0.36

    fig, ax1 = plt.subplots(figsize=(6.6, 3.4), dpi=130)
    fig.patch.set_facecolor("#EFEFEF")
    ax1.set_facecolor("#EFEFEF")

    gross_bars = ax1.bar(
        [i - (width / 2.0) for i in x],
        gross,
        width=width,
        color=f"#{ACCENT_HEX}",
        alpha=0.9,
        label=primary_label,
        zorder=2,
    )

    ax2 = ax1.twinx()
    line_x = [i + (width / 2.0) + 0.06 for i in x]
    ax2.plot(
        line_x,
        quantity,
        color=f"#{SECONDARY_HEX}",
        marker="o",
        linewidth=2.1,
        markersize=4.5,
        label="Quantity",
        zorder=3,
    )

    ax1.set_xticks(x)
    rotate_x = 16 if max((len(str(label or "")) for label in labels), default=0) > 12 else 0
    ax1.set_xticklabels(
        [_short_label(label, max_len=14) for label in labels],
        fontsize=8,
        rotation=rotate_x,
        ha="right" if rotate_x else "center",
    )
    ax1.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_money_tick(value)))
    ax2.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_quantity_tick(value)))
    ax1.set_ylabel(primary_label, fontsize=8, color="#4A586F")
    ax2.set_ylabel("Quantity", fontsize=8, color="#4A586F")
    ax1.tick_params(axis="y", labelsize=8, colors="#4A586F")
    ax2.tick_params(axis="y", labelsize=8, colors="#4A586F")
    ax1.grid(axis="y", linestyle="--", linewidth=0.6, alpha=0.28, zorder=1)
    ax1.set_title(title, fontsize=10, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=6)
    gross_max = max(gross, default=0.0)
    qty_max = max(quantity, default=0.0)
    ax1.set_ylim(0, gross_max * 1.22 if gross_max > 0 else 1.0)
    ax2.set_ylim(0, qty_max * 1.4 if qty_max > 0 else 1.0)
    ax1.set_xlim(-0.5, (len(labels) - 0.5) + 0.24)

    for side in ("top", "right"):
        ax1.spines[side].set_visible(False)
    ax1.spines["left"].set_alpha(0.3)
    ax1.spines["bottom"].set_alpha(0.3)
    ax2.spines["top"].set_visible(False)
    ax2.spines["left"].set_visible(False)
    ax2.spines["right"].set_alpha(0.3)

    legend_handles = [
        Patch(facecolor=f"#{ACCENT_HEX}", edgecolor=f"#{ACCENT_HEX}", label=primary_label),
        Line2D([0], [0], color=f"#{SECONDARY_HEX}", marker="o", linewidth=2.1, markersize=4.5, label="Quantity"),
    ]
    ax1.legend(
        handles=legend_handles,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.18),
        frameon=False,
        ncol=2,
        fontsize=8.5,
        handlelength=1.2,
        columnspacing=0.9,
    )

    for bar in gross_bars:
        value = float(bar.get_height() or 0.0)
        if value <= 0:
            continue
        x_pos = bar.get_x() + bar.get_width() / 2.0
        ax1.text(
            x_pos,
            value,
            _compact_number(value),
            ha="center",
            va="bottom",
            fontsize=7,
            color="#374151",
        )

    for idx, value in enumerate(quantity):
        if value <= 0:
            continue
        label_y = value + (qty_max * 0.04 if qty_max > 0 else 0.0)
        label_y = min(label_y, ax2.get_ylim()[1] * 0.98)
        ax2.text(
            line_x[idx],
            label_y,
            _format_quantity_number(value),
            ha="center",
            va="bottom",
            fontsize=7,
            color=f"#{SECONDARY_HEX}",
            bbox={"boxstyle": "round,pad=0.12", "facecolor": "#EFEFEF", "edgecolor": "none", "alpha": 0.95},
        )

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_pie_chart_image(rows: list[dict[str, Any]], title: str) -> BytesIO:
    ordered = sorted(rows, key=lambda item: _to_float(item.get("gross_premium")), reverse=True)[:6]
    labels = [str(item.get("label") or "") for item in ordered if str(item.get("label") or "").strip()]
    values = [_to_float(item.get("gross_premium")) for item in ordered[: len(labels)]]

    if not labels or sum(values) <= 0:
        labels = ["No data"]
        values = [1.0]

    colors = ["#2D6BE8", "#F97316", "#FBBF24", "#22C55E", "#14B8A6", "#A855F7"]
    fig, ax = plt.subplots(figsize=(5.4, 3.2), dpi=130)
    fig.patch.set_facecolor("#EFEFEF")
    ax.set_facecolor("#EFEFEF")

    wedges, _, _ = ax.pie(
        values,
        startangle=90,
        autopct=lambda pct: f"{pct:.0f}%",
        colors=colors[: len(values)],
        textprops={"fontsize": 8},
        wedgeprops={"linewidth": 0.8, "edgecolor": "#EFEFEF"},
    )
    ax.legend(
        wedges,
        [_short_label(label, max_len=16) for label in labels],
        loc="center left",
        bbox_to_anchor=(1.02, 0.5),
        fontsize=8,
        frameon=False,
    )
    ax.set_title(title, fontsize=10, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=6)
    ax.axis("equal")

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_radial_chart_image(rows: list[dict[str, Any]], title: str) -> BytesIO:
    ordered = sorted(rows, key=lambda item: _to_float(item.get("gross_premium")), reverse=True)[:6]
    labels = [str(item.get("label") or "") for item in ordered if str(item.get("label") or "").strip()]
    values = [_to_float(item.get("gross_premium")) for item in ordered[: len(labels)]]

    if not labels or max(values, default=0.0) <= 0:
        labels = ["No data"]
        values = [1.0]

    count = max(1, len(values))
    angles = [(2.0 * math.pi * idx) / count for idx in range(count)]
    max_value = max(values) if values else 1.0
    normalized = [(value / max_value) if max_value else 0.0 for value in values]

    fig, ax = plt.subplots(figsize=(5.4, 3.2), dpi=130, subplot_kw={"projection": "polar"})
    fig.patch.set_facecolor("#EFEFEF")
    ax.set_facecolor("#EFEFEF")
    ax.set_theta_offset(math.pi / 2.0)
    ax.set_theta_direction(-1)
    bar_width = (2.0 * math.pi / count) * 0.72

    bars = ax.bar(angles, normalized, width=bar_width, color=f"#{ACCENT_HEX}", alpha=0.86)
    ax.set_xticks(angles)
    ax.set_xticklabels([_short_label(label, max_len=12) for label in labels], fontsize=8, color="#374151")
    ax.set_yticklabels([])
    ax.grid(alpha=0.22)
    ax.spines["polar"].set_alpha(0.2)
    ax.set_title(title, fontsize=10, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=8)

    for bar, value in zip(bars, values):
        angle = bar.get_x() + bar.get_width() / 2.0
        radius = min(1.2, bar.get_height() + 0.11)
        ax.text(angle, radius, _compact_number(value), ha="center", va="center", fontsize=7, color="#374151")

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_asp_chart_image(rows: list[dict[str, Any]], *, dimension: str, title: str) -> BytesIO:
    labels = [str(row.get("label") or "") for row in rows]
    asp_values = [_to_float(row.get("asp")) for row in rows]

    if not labels:
        labels = ["No data"]
        asp_values = [0.0]

    fig, ax = plt.subplots(figsize=(8.4, 2.65), dpi=130)
    fig.patch.set_facecolor("#F5F7FA")
    ax.set_facecolor("#F5F7FA")

    x = list(range(len(labels)))
    max_asp = max(asp_values, default=0.0)
    zero_floor = max(max_asp * 0.015, 0.2) if max_asp > 0 else 0.2
    if dimension in {"month", "week"}:
        ax.plot(x, asp_values, color=f"#{SECONDARY_HEX}", marker="o", linewidth=2.2, markersize=4.5, zorder=3)
        ax.fill_between(x, asp_values, color=f"#{SECONDARY_HEX}", alpha=0.12, zorder=2)
    else:
        display_values = [value if value > 0 else zero_floor for value in asp_values]
        bars = ax.bar(x, display_values, color=f"#{SECONDARY_HEX}", width=0.62, zorder=2)
        for idx, bar in enumerate(bars):
            if asp_values[idx] > 0:
                continue
            bar.set_alpha(0.35)

    ax.set_xticks(x)
    ax.set_xticklabels([_short_label(label, max_len=14) for label in labels], fontsize=9, color="#4A5B74")
    ax.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_money_tick(value)))
    ax.set_ylabel("ASP", fontsize=8.5, color="#4A5B74")
    ax.tick_params(axis="y", labelsize=9, colors="#4A5B74")
    ax.grid(axis="y", linestyle=(0, (3, 3)), linewidth=0.7, alpha=0.28, zorder=1)
    ax.set_title(title, fontsize=10, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=6)
    y_max = max(max(asp_values, default=0.0), zero_floor) * 1.2
    ax.set_ylim(0, y_max if y_max > 0 else 1.0)

    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.spines["left"].set_alpha(0.3)
    ax.spines["bottom"].set_alpha(0.3)

    for idx, value in enumerate(asp_values):
        label_value = value if value > 0 else 0.0
        label_y = value if value > 0 else (max(max_asp * 0.04, 0.15))
        ax.text(
            x[idx],
            label_y,
            _compact_number(label_value),
            ha="center",
            va="bottom",
            fontsize=7,
            color=f"#{SECONDARY_HEX}",
        )

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_metric_bar_chart_image(
    rows: list[dict[str, Any]],
    *,
    metric_key: str,
    title: str,
    color_hex: str,
    quantity_mode: bool,
    primary_label: str = "Gross Premium",
) -> BytesIO:
    labels = [str(row.get("label") or "") for row in rows]
    values = [_to_float(row.get(metric_key)) for row in rows]
    if not labels:
        labels = ["No data"]
        values = [0.0]

    fig, ax = plt.subplots(figsize=(5.95, 3.1), dpi=130)
    fig.patch.set_facecolor("#F5F7FA")
    ax.set_facecolor("#F5F7FA")
    x = list(range(len(labels)))
    bars = ax.bar(x, values, color=f"#{color_hex}", width=0.62, alpha=0.88, zorder=2)
    ax.set_xticks(x)
    ax.set_xticklabels([_short_label(label, max_len=12) for label in labels], fontsize=8.2, color="#4A5B74")
    ax.grid(axis="y", linestyle=(0, (3, 3)), linewidth=0.7, alpha=0.28, zorder=1)
    if quantity_mode:
        ax.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_quantity_tick(value)))
        ax.set_ylabel("No. of Plans", fontsize=8, color="#4A5B74")
    else:
        ax.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_money_tick(value)))
        ax.set_ylabel(primary_label, fontsize=8, color="#4A5B74")
    ax.tick_params(axis="y", labelsize=8.2, colors="#4A5B74")
    ax.set_title(title, fontsize=9.8, fontweight="bold", color=f"#{TITLE_HEX}", loc="left", pad=6)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.spines["left"].set_alpha(0.3)
    ax.spines["bottom"].set_alpha(0.3)
    max_value = max(values, default=0.0)
    ax.set_ylim(0, max_value * 1.2 if max_value > 0 else 1.0)

    for bar in bars:
        value = float(bar.get_height() or 0.0)
        if value <= 0:
            continue
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            value,
            _format_quantity_number(value) if quantity_mode else _compact_number(value),
            ha="center",
            va="bottom",
            fontsize=6.9,
            color="#334155",
        )

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _render_combo_chart_image(rows: list[dict[str, Any]], *, primary_label: str = "Gross Premium") -> BytesIO:
    labels = [str(row.get("label") or "") for row in rows]
    gross = [_to_float(row.get("gross_premium")) for row in rows]
    quantity = [_to_float(row.get("quantity")) for row in rows]

    if not labels:
        labels = ["No data"]
        gross = [0.0]
        quantity = [0.0]

    x = list(range(len(labels)))
    fig, ax1 = plt.subplots(figsize=(8.4, 2.65), dpi=130)
    fig.patch.set_facecolor("#F5F7FA")
    ax1.set_facecolor("#F5F7FA")

    bars = ax1.bar(x, gross, color=f"#{ACCENT_HEX}", width=0.64, label=primary_label, zorder=2)
    ax1.grid(axis="y", linestyle=(0, (3, 3)), linewidth=0.7, alpha=0.28, zorder=1)
    ax1.set_xticks(x)
    rotate_x = 16 if max((len(str(label or "")) for label in labels), default=0) > 12 else 0
    ax1.set_xticklabels(
        [_short_label(label, max_len=14) for label in labels],
        fontsize=9,
        color="#4A5B74",
        rotation=rotate_x,
        ha="right" if rotate_x else "center",
    )
    ax1.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_money_tick(value)))
    ax1.set_ylabel(primary_label, fontsize=8.5, color="#4A5B74")
    ax1.tick_params(axis="y", labelsize=9, colors="#4A5B74")
    for side in ("top", "right"):
        ax1.spines[side].set_visible(False)
    ax1.spines["left"].set_alpha(0.3)
    ax1.spines["bottom"].set_alpha(0.3)

    ax2 = ax1.twinx()
    line_x = [i + 0.18 for i in x]
    ax2.plot(line_x, quantity, color=f"#{SECONDARY_HEX}", marker="o", markersize=4.2, linewidth=2.0, label="Quantity", zorder=3)
    ax2.yaxis.set_major_formatter(FuncFormatter(lambda value, _: _axis_quantity_tick(value)))
    ax2.set_ylabel("Quantity", fontsize=8.5, color="#4A5B74")
    ax2.tick_params(axis="y", labelsize=9, colors="#4A5B74")
    ax2.spines["top"].set_visible(False)
    ax2.spines["left"].set_visible(False)
    ax2.spines["right"].set_alpha(0.3)
    gross_max = max(gross, default=0.0)
    qty_max = max(quantity, default=0.0)
    ax1.set_ylim(0, gross_max * 1.2 if gross_max > 0 else 1.0)
    ax2.set_ylim(0, qty_max * 1.4 if qty_max > 0 else 1.0)
    ax1.set_xlim(-0.5, (len(labels) - 0.5) + 0.24)

    legend_handles = [
        Patch(facecolor=f"#{ACCENT_HEX}", edgecolor=f"#{ACCENT_HEX}", label=primary_label),
        Line2D([0], [0], color=f"#{SECONDARY_HEX}", marker="o", linewidth=2.0, markersize=4.2, label="Quantity"),
    ]
    ax1.legend(
        handles=legend_handles,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.16),
        frameon=False,
        ncol=2,
        fontsize=9.5,
        handlelength=1.4,
        columnspacing=1.0,
    )

    for bar in bars:
        value = float(bar.get_height() or 0.0)
        if value <= 0:
            continue
        ax1.text(
            bar.get_x() + bar.get_width() / 2.0,
            value,
            _compact_number(value),
            ha="center",
            va="bottom",
            fontsize=7.2,
            color="#334155",
        )

    for idx, value in enumerate(quantity):
        if value <= 0:
            continue
        label_y = value + (qty_max * 0.04 if qty_max > 0 else 0.0)
        label_y = min(label_y, ax2.get_ylim()[1] * 0.98)
        ax2.text(
            line_x[idx],
            label_y,
            _format_quantity_number(value),
            ha="center",
            va="bottom",
            fontsize=7.2,
            color=f"#{SECONDARY_HEX}",
            bbox={"boxstyle": "round,pad=0.12", "facecolor": "#F5F7FA", "edgecolor": "none", "alpha": 0.95},
        )

    image = BytesIO()
    plt.tight_layout()
    fig.savefig(image, format="png", dpi=130, facecolor=fig.get_facecolor())
    plt.close(fig)
    image.seek(0)
    return image


def _build_asp_insights(rows: list[dict[str, Any]], dimension: str) -> list[str]:
    points = [(str(row.get("label") or ""), _to_float(row.get("asp"))) for row in rows]
    points = [(label, value) for label, value in points if label]
    if not points:
        return ["No ASP points were found for the selected period."]

    insights: list[str] = []
    if dimension in {"month", "week"} and len(points) >= 2:
        latest_label, latest = points[-1]
        prev_label, prev = points[-2]
        delta = latest - prev
        pct = (delta / prev * 100.0) if prev else 0.0
        direction = "up" if delta >= 0 else "down"
        insights.append(
            f"In the latest period ({latest_label}), ASP moved {direction} by {_compact_number(abs(delta))} ({pct:+.1f}%) versus {prev_label}, indicating the most recent pricing direction."
        )

        first_label, first = points[0]
        overall_dir = "increasing" if latest >= first else "softening"
        insights.append(
            f"Across the selected date range, ASP is {overall_dir}, moving from {first_label} to {latest_label}."
        )

    max_label, max_value = max(points, key=lambda item: item[1])
    min_label, min_value = min(points, key=lambda item: item[1])
    insights.append(
        f"The highest ASP bucket is {max_label} at {_compact_number(max_value)}, while {min_label} is lowest at {_compact_number(min_value)}, showing the spread across categories."
    )

    values = [value for _, value in points]
    avg = sum(values) / len(values)
    spread = max(values) - min(values)
    insights.append(
        f"Average ASP stands at {_compact_number(avg)}, with an overall spread of {_compact_number(spread)} between the highest and lowest buckets."
    )

    ordered = sorted(points, key=lambda item: item[1], reverse=True)
    top = ", ".join(label for label, _ in ordered[:3])
    bottom = ", ".join(label for label, _ in ordered[-3:]) if len(ordered) >= 3 else top
    insights.append(f"Top ASP contributors are {top}.")
    insights.append(f"Lower ASP contributors are {bottom}.")
    return insights[:7]


def _build_trend_insights(
    rows: list[dict[str, Any]],
    dimension: str,
    metric: str,
    *,
    metric_label: str | None = None,
) -> list[str]:
    points = _to_points(rows, dimension, metric)
    if not points:
        return ["No comparable trend points were found for the selected drilldown."]

    display_metric = (metric_label or _prettify_label(metric)).strip() or _prettify_label(metric)
    insights: list[str] = []
    if dimension in {"month", "week"} and len(points) >= 2:
        latest_label, latest_value = points[-1]
        prev_label, prev_value = points[-2]
        delta = latest_value - prev_value
        pct = (delta / prev_value * 100.0) if prev_value else 0.0
        direction = "up" if delta >= 0 else "down"
        insights.append(
            f"In the latest period ({latest_label}), {display_metric} moved {direction} by {_compact_number(abs(delta))} ({pct:+.1f}%) versus {prev_label}, which defines the immediate trend shift."
        )

        first_label, first_value = points[0]
        overall_delta = latest_value - first_value
        overall_dir = "increasing" if overall_delta >= 0 else "softening"
        insights.append(
            f"Across the full selected timeline, {display_metric} is {overall_dir}, shifting from {first_label} to {latest_label}."
        )
    else:
        max_label, max_value = max(points, key=lambda item: item[1])
        min_label, min_value = min(points, key=lambda item: item[1])
        insights.append(
            f"The leading drilldown is {max_label} at {_compact_number(max_value)}, while {min_label} remains at the lower end with {_compact_number(min_value)}."
        )
        if min_value > 0:
            ratio = max_value / min_value
            insights.append(
                f"The top bucket contributes about {ratio:.1f}x the lowest bucket, indicating visible concentration in a few high-performing segments."
            )
        else:
            insights.append("Contribution is highly skewed, with top segments carrying most of the selected metric.")

    sorted_points = sorted(points, key=lambda item: item[1], reverse=True)
    top_three = ", ".join(str(label) for label, _ in sorted_points[:3])
    bottom_three = ", ".join(str(label) for label, _ in sorted_points[-3:]) if len(sorted_points) >= 3 else top_three
    insights.append(f"The top contributors are {top_three}.")
    insights.append(f"The comparatively weaker contributors are {bottom_three}.")

    values = [v for _, v in points]
    avg = sum(values) / len(values)
    volatility = (max(values) - min(values)) if values else 0.0
    insights.append(
        f"Average {display_metric} is {_compact_number(avg)}, with a spread of {_compact_number(volatility)} from peak to trough."
    )

    if len(values) >= 3:
        top_total = sum(sorted(values, reverse=True)[:3])
        total = sum(values) or 1.0
        share = top_total / total * 100.0
        insights.append(f"The top 3 buckets together contribute {share:.1f}% of total {display_metric}.")

        median = float(pd.Series(values).median())
        insights.append(f"The median bucket value for {display_metric} is {_compact_number(median)}.")

    if dimension in {"month", "week"} and len(values) >= 4:
        lead = sum(values[-2:]) / 2.0
        base = sum(values[:2]) / 2.0
        change = ((lead - base) / base * 100.0) if base else 0.0
        insights.append(
            f"Momentum check: recent periods are {change:+.1f}% versus the earliest baseline periods in the selected range."
        )

    return insights[:7]


def _to_points(rows: list[dict[str, Any]], dimension: str, metric: str) -> list[tuple[str, float]]:
    out: list[tuple[str, float]] = []
    dim_key = _safe_key(dimension)
    metric_key = _safe_key(metric)

    for row in rows:
        if not isinstance(row, dict):
            continue
        safe_map = {_safe_key(str(k)): k for k in row.keys()}
        dim_candidates = [dim_key]
        if dim_key == "model_code":
            dim_candidates.extend(["product_category", "product_description", "plan_product", "model"])
        dim_col = (
            next((safe_map.get(candidate) for candidate in dim_candidates if safe_map.get(candidate)), None)
            or safe_map.get("week")
            or safe_map.get("month")
            or safe_map.get("state")
            or safe_map.get("label")
        )
        metric_col = safe_map.get(metric_key) or safe_map.get("value") or safe_map.get("quantity")
        if dim_col is None or metric_col is None:
            continue
        label = str(row.get(dim_col) or "").strip()
        if not label:
            continue
        value = _to_float(row.get(metric_col))
        out.append((label, value))

    if dimension == "month":
        out.sort(key=lambda item: _month_sort_key(item[0]))
    elif dimension == "week":
        out.sort(key=lambda item: _week_sort_key(item[0]))
    return out


def _canonical_plan_category_label(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "").strip().lower())
    if not text:
        return ""
    if "combo" in text:
        return "COMBO"
    if "adld" in text or "accidental" in text or "liquid" in text:
        return "ADLD"
    if re.search(r"\bsp\b|\bspp\b", text) or "screen" in text or "crack" in text:
        return "SP"
    if re.search(r"\bew\b", text) or "extended warranty" in text or text.startswith("ew"):
        return "EW"
    return str(value or "").strip()


def _canonical_device_plan_category_label(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "").strip().lower())
    if not text:
        return ""
    if "luxury" in text and "fold" in text:
        return "Luxury Fold"
    if "luxury" in text and "flip" in text:
        return "Luxury Flip"
    if "fold" in text:
        return "Luxury Fold"
    if "flip" in text:
        return "Luxury Flip"
    if "super" in text and "premium" in text:
        return "Super Premium"
    if text.startswith("premium") or ("premium" in text and "super" not in text):
        return "Premium"
    if text.startswith("high") or text == "high":
        return "High"
    if text.startswith("mid") or text == "mid":
        return "Mid"
    if text.startswith("mass") or text == "mass":
        return "Mass"
    return str(value or "").strip()


def _canonical_dimension_label(label: str, dimension: str | None) -> str:
    dim_key = _safe_key(str(dimension or ""))
    raw = str(label or "").strip()
    if not raw:
        return ""
    if dim_key == "plan_category":
        return _canonical_plan_category_label(raw)
    if dim_key == "device_plan_category":
        return _canonical_device_plan_category_label(raw)
    return raw


def _merge_metric_points(
    gross_points: list[tuple[str, float]],
    quantity_points: list[tuple[str, float]],
    *,
    dimension: str | None = None,
) -> list[dict[str, Any]]:
    merged: dict[str, dict[str, Any]] = {}
    order: list[str] = []

    for label, value in gross_points:
        key = _canonical_dimension_label(str(label), dimension)
        if not key:
            continue
        if key not in merged:
            merged[key] = {"label": key, "gross_premium": 0.0, "quantity": 0.0}
            order.append(key)
        merged[key]["gross_premium"] = _to_float(value)

    for label, value in quantity_points:
        key = _canonical_dimension_label(str(label), dimension)
        if not key:
            continue
        if key not in merged:
            merged[key] = {"label": key, "gross_premium": 0.0, "quantity": 0.0}
            order.append(key)
        merged[key]["quantity"] = _to_float(value)

    return [merged[key] for key in order]


def _safe_key(value: str) -> str:
    return re.sub(r"[()%'.]", "", re.sub(r"\s+", "_", (value or "").strip().lower()))


def _device_plan_category_rank(label: str) -> int:
    normalized = re.sub(r"\s+", " ", str(label or "").strip().lower())
    for idx, expected in enumerate(DEVICE_PLAN_CATEGORY_ORDER):
        if normalized == expected.lower():
            return idx
    return 99


def _to_float(value: Any) -> float:
    try:
        n = float(value or 0.0)
    except Exception:
        return 0.0
    if not math.isfinite(n):
        return 0.0
    return n


def _short_label(value: str, max_len: int = 16) -> str:
    text = str(value or "").strip()
    return text if len(text) <= max_len else f"{text[: max_len - 1]}..."


def _prettify_label(value: str) -> str:
    label = str(value or "").replace("_", " ").strip().title()
    if _safe_key(value) in {"model_code", "model_code_1", "model", "product_model"}:
        return "Product Model"
    return label


def _compact_number(value: float) -> str:
    n = float(value or 0.0)
    abs_n = abs(n)
    if abs_n >= 1e7:
        return f"{n/1e7:.2f} Cr"
    if abs_n >= 1e5:
        return f"{n/1e5:.2f} L"
    return f"{n:,.0f}"


def _format_quantity_number(value: float) -> str:
    n = float(value or 0.0)
    abs_n = abs(n)
    if abs_n >= 1e7:
        return f"{n/1e7:.2f} Cr"
    if abs_n >= 1e5:
        return f"{n/1e5:.2f} L"
    return f"{int(round(n)):,}"


def _format_metric_value(label: str, value: float) -> str:
    low = label.lower()
    if "unit" in low or "quantity" in low:
        return _format_quantity_number(value)
    return f"Rs {_compact_number(value)}"


def _axis_money_tick(value: float) -> str:
    return _compact_number(_to_float(value))


def _axis_quantity_tick(value: float) -> str:
    return _format_quantity_number(_to_float(value))


def _month_sort_key(value: str) -> tuple[int, int, str]:
    raw = str(value or "").strip()
    if not raw:
        return (9999, 12, raw)

    ts = pd.to_datetime(raw, errors="coerce")
    if pd.notna(ts):
        return (int(ts.year), int(ts.month), raw)

    m = re.match(r"^([A-Za-z]{3})[-/ ](\d{2,4})$", raw)
    if m:
        month_map = {
            "jan": 1,
            "feb": 2,
            "mar": 3,
            "apr": 4,
            "may": 5,
            "jun": 6,
            "jul": 7,
            "aug": 8,
            "sep": 9,
            "oct": 10,
            "nov": 11,
            "dec": 12,
        }
        month = month_map.get(m.group(1).lower(), 12)
        year_token = int(m.group(2))
        year = year_token + 2000 if year_token < 100 else year_token
        return (year, month, raw)

    return (9999, 12, raw)


def _top_label(rows: list[dict[str, Any]], dimension: str) -> str | None:
    points = _to_points(rows, dimension, "quantity") or _to_points(rows, dimension, "gross_premium")
    if not points:
        return None
    points_sorted = sorted(points, key=lambda item: item[1], reverse=True)
    return points_sorted[0][0] if points_sorted else None
